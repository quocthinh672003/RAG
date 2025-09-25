# POST /export
import os
import uuid
import json
import csv
import io
from fastapi import APIRouter
from app.core.schemas import ExportIn
from app.core.openai_client import client

router = APIRouter(prefix="/export", tags=["export"])


async def _download_openai_file(file_id: str) -> bytes:
    """Download file from OpenAI"""
    stream = await client.files.content(file_id)
    if hasattr(stream, "aread"):
        return await stream.aread()
    if hasattr(stream, "read"):
        return stream.read()
    if isinstance(stream, (bytes, bytearray)):
        return bytes(stream)
    raise RuntimeError("Unsupported file content stream type")


def _local_generate(fmt: str, content: str, path: str) -> None:
    """Generate file locally as fallback"""
    os.makedirs(os.path.dirname(path) or ".", exist_ok=True)
    
    if fmt == "md":
        with open(path, "w", encoding="utf-8") as f:
            f.write(content or "")
    
    elif fmt == "html":
        html_doc = f"""<!doctype html><html><head><meta charset='utf-8'><title>Export</title></head><body>{content or ""}</body></html>"""
        with open(path, "w", encoding="utf-8") as f:
            f.write(html_doc)
    
    elif fmt == "csv":
        with open(path, "w", newline="", encoding="utf-8") as fh:
            writer = csv.writer(fh)
            for line in io.StringIO(content or "").read().strip().splitlines():
                writer.writerow([p.strip() for p in line.split(',')])
    
    elif fmt == "xlsx":
        import pandas as pd
        rows = []
        for line in io.StringIO(content or "").read().strip().splitlines():
            rows.append([p.strip() for p in line.split(',')])
        df = pd.DataFrame(rows[1:], columns=rows[0] if rows else None)
        df.to_excel(path, index=False)
    
    elif fmt == "pdf":
        try:
            from reportlab.lib.pagesizes import A4
            from reportlab.pdfgen import canvas
            from reportlab.lib.units import cm
            c = canvas.Canvas(path, pagesize=A4)
            w, h = A4
            t = c.beginText(2*cm, h-2*cm)
            for line in (content or "").splitlines():
                t.textLine(line)
            c.drawText(t)
            c.showPage()
            c.save()
        except Exception:
            pass
    
    elif fmt == "docx":
        from docx import Document
        doc = Document()
        for line in (content or "").splitlines():
            doc.add_paragraph(line)
        doc.save(path)
    
    elif fmt == "pptx":
        from pptx import Presentation
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = ((content or "").splitlines() or ["Export"])[0]
        tf = slide.placeholders[1].text_frame
        for l in (content or "").splitlines()[1:]:
            tf.add_paragraph().text = l
        prs.save(path)


def _collect_output_files(response) -> list:
    """Collect file artifacts from response"""
    files = []
    for block in (getattr(response, "output", []) or []):
        for c in (getattr(block, "content", []) or []):
            f = getattr(c, "file", None)
            if f:
                files.append({
                    "id": getattr(f, "id", None),
                    "name": getattr(f, "filename", None),
                })
    return files


@router.post("/")
async def export_content(payload: ExportIn):
    """Export content to various formats"""
    
    os.makedirs("storage", exist_ok=True)
    export_id = uuid.uuid4().hex[:8]
    fmt = payload.format.lower()
    
    filename = f"export_{export_id}.{fmt}"
    path = os.path.join("storage", filename)
    src = json.dumps(payload.content or "")
    
    # Code templates for Code Interpreter
    templates = {
        "html": f"""
SRC = {src}
html_doc = "<!doctype html><html><head><meta charset='utf-8'><title>Export</title></head><body>" + SRC + "</body></html>"
with open('{filename}', 'w', encoding='utf-8') as f:
    f.write(html_doc)
""",
        "csv": f"""
import csv
SRC = {src}
with open('{filename}', 'w', newline='', encoding='utf-8') as f:
    w = csv.writer(f)
    for line in (SRC or '').strip().splitlines():
        w.writerow([p.strip() for p in line.split(',')])
""",
        "xlsx": f"""
import pandas as pd
SRC = {src}
rows = [[p.strip() for p in line.split(',')] for line in (SRC or '').strip().splitlines()]
df = pd.DataFrame(rows[1:], columns=rows[0] if rows else None)
df.to_excel('{filename}', index=False)
""",
        "pdf": f"""
from reportlab.lib.pagesizes import A4
from reportlab.pdfgen import canvas
from reportlab.lib.units import cm
SRC = {src}
c = canvas.Canvas('{filename}', pagesize=A4)
w, h = A4
t = c.beginText(2*cm, h-2*cm)
for line in (SRC or '').splitlines():
    t.textLine(line)
c.drawText(t); c.showPage(); c.save()
""",
        "docx": f"""
from docx import Document
SRC = {src}
doc = Document()
for line in (SRC or '').splitlines():
    doc.add_paragraph(line)
doc.save('{filename}')
""",
        "pptx": f"""
from pptx import Presentation
SRC = {src}
prs = Presentation()
s = prs.slides.add_slide(prs.slide_layouts[1])
s.shapes.title.text = ((SRC or '').splitlines() or ['Export'])[0]
tf = s.placeholders[1].text_frame
for l in (SRC or '').splitlines()[1:]:
    tf.add_paragraph().text = l
prs.save('{filename}')
"""
    }
    
    code_block = templates.get(fmt, f"""
SRC = {src}
with open('{filename}', 'w', encoding='utf-8') as f:
    f.write(SRC)
""")
    
    instruction = f"""
Execute this Python code in Code Interpreter.
Save the file exactly as {filename}.
Return it as a file artifact (not text, not markdown link).

```python
{code_block}
```
""".strip()
    
    artifacts = []
    response = None
    
    try:
        response = await client.responses.create(
            model="gpt-4.1",
            input=instruction,
            tools=[{"type": "code_interpreter", "container": {"image": "default"}}],
            tool_choice="code_interpreter",
            temperature=0,
        )
        artifacts = _collect_output_files(response)
    except Exception:
        artifacts = []
    
    # If Code Interpreter returned a file, download and return it
    if artifacts:
        picked = next((a for a in artifacts if (a["name"] or "").lower().endswith(f".{fmt}")), artifacts[0])
        blob = await _download_openai_file(picked["id"])
        with open(path, "wb") as fh:
            fh.write(blob)
        
        # PDF integrity check
        if fmt == "pdf":
            try:
                size_ok = os.path.exists(path) and os.path.getsize(path) > 100
                head_ok = False
                if os.path.exists(path):
                    with open(path, "rb") as fh:
                        head_ok = (fh.read(5) == b"%PDF-")
                if not (size_ok and head_ok):
                    from reportlab.lib.pagesizes import A4
                    from reportlab.pdfgen import canvas
                    from reportlab.lib.units import cm
                    c = canvas.Canvas(path, pagesize=A4)
                    w, h = A4
                    t = c.beginText(2*cm, h-2*cm)
                    for line in (payload.content or "").splitlines():
                        t.textLine(line)
                    c.drawText(t)
                    c.showPage()
                    c.save()
            except Exception:
                pass
        
        return {
            "download_url": f"/static/{os.path.basename(path)}",
            "response_id": getattr(response, "id", None),
            "filename": os.path.basename(path),
        }
    
    # Fallback: generate locally
    _local_generate(fmt, payload.content or "", path)
    return {
        "download_url": f"/static/{os.path.basename(path)}",
        "response_id": getattr(response, "id", None),
        "filename": os.path.basename(path),
    }